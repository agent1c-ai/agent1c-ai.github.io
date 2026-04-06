from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "hackathon" / "agent1c-solana-hackathon-pitch-deck.pptx"
HEDGEHOG_IMAGE = ROOT / "assets" / "hedgey-clippy.png"


def rgb(value):
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BG_DARK = rgb("0B1020")
BG_DARK_2 = rgb("121B31")
BG_LIGHT = rgb("F7F3EA")
BG_SOFT = rgb("EAF2F5")
INK_LIGHT = rgb("F7F4EE")
INK_DARK = rgb("14233B")
INK_MUTED = rgb("4F637D")
MINT = rgb("35E4C4")
SKY = rgb("69CFFF")
CORAL = rgb("FF7A59")
GOLD = rgb("FFC857")
GREEN = rgb("5ED27A")
PANEL_DARK = rgb("17243D")
PANEL_LIGHT = rgb("FFFFFF")
PANEL_SOFT = rgb("EEF4F7")

TITLE_FONT = "Trebuchet MS"
BODY_FONT = "Trebuchet MS"


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.18
    return shape


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    kind = SHAPE.ROUNDED_RECTANGLE if radius else SHAPE.RECTANGLE
    return add_shape(slide, kind, x, y, w, h, fill, line=line, radius=radius)


def add_oval(slide, x, y, w, h, fill, line=None):
    return add_shape(slide, SHAPE.OVAL, x, y, w, h, fill, line=line)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK_DARK,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    lines = str(text).split("\n")
    for idx, line in enumerate(lines):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = line
        para.alignment = align
        para.space_after = Pt(0)
        para.line_spacing = line_spacing
        font_obj = para.font
        font_obj.name = font
        font_obj.size = Pt(size)
        font_obj.color.rgb = color
        font_obj.bold = bold
    return box


def add_bullets(slide, items, x, y, w, h, *, size=18, color=INK_DARK, leading=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for idx, item in enumerate(items):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = f"• {item}"
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(8)
        para.line_spacing = leading
        font_obj = para.font
        font_obj.name = BODY_FONT
        font_obj.size = Pt(size)
        font_obj.color.rgb = color
    return box


def add_kicker(slide, text, x, y, w=2.25, h=0.38, fill=MINT, text_color=BG_DARK):
    add_rect(slide, x, y, w, h, fill, radius=True)
    add_text(
        slide,
        text,
        x + 0.13,
        y + 0.06,
        w - 0.26,
        h - 0.12,
        size=11,
        color=text_color,
        bold=True,
        font=BODY_FONT,
    )


def add_footer(slide, page, *, appendix=False, dark=False):
    color = INK_LIGHT if dark else INK_MUTED
    label = "Appendix" if appendix else "Pitch Deck"
    add_text(
        slide,
        f"Agent1c.ai  |  Solana Hackathon  |  {label}",
        0.7,
        7.05,
        5.0,
        0.22,
        size=10,
        color=color,
    )
    add_text(
        slide,
        str(page),
        12.15,
        7.02,
        0.4,
        0.24,
        size=10,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def decorate_dark(slide):
    add_rect(slide, 0, 0, 13.333, 7.5, BG_DARK)
    add_oval(slide, 10.4, -1.1, 4.2, 4.2, BG_DARK_2)
    add_oval(slide, -0.9, 5.65, 2.4, 2.4, rgb("16223A"))
    add_rect(slide, 0.7, 0.58, 0.65, 0.08, MINT)
    add_rect(slide, 1.42, 0.58, 0.46, 0.08, CORAL)
    add_rect(slide, 1.94, 0.58, 0.36, 0.08, GOLD)


def decorate_light(slide, alt=False):
    add_rect(slide, 0, 0, 13.333, 7.5, BG_SOFT if alt else BG_LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.18, MINT if alt else CORAL)
    add_oval(slide, 11.25, -0.65, 2.4, 2.4, rgb("E4ECEE") if alt else rgb("EFE6DB"))
    add_oval(slide, -0.55, 6.18, 1.8, 1.8, rgb("DDEBED") if alt else rgb("F0E0D1"))


def add_title_block(slide, title, subtitle, *, dark=False):
    title_color = INK_LIGHT if dark else INK_DARK
    sub_color = rgb("B9C8DB") if dark else INK_MUTED
    add_text(slide, title, 0.7, 1.0, 7.0, 1.2, size=30, color=title_color, bold=True, font=TITLE_FONT)
    add_text(slide, subtitle, 0.7, 1.95, 6.2, 0.9, size=14, color=sub_color)


def add_card(slide, x, y, w, h, title, body, *, fill=PANEL_LIGHT, title_color=INK_DARK, body_color=INK_MUTED):
    add_rect(slide, x, y, w, h, fill, radius=True)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.38, size=14, color=title_color, bold=True)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.74, size=12, color=body_color)


def add_number_card(slide, num, title, body, x, y, w, h, *, fill=PANEL_DARK, accent=MINT):
    add_rect(slide, x, y, w, h, fill, radius=True)
    add_oval(slide, x + 0.18, y + 0.18, 0.52, 0.52, accent)
    add_text(slide, str(num), x + 0.18, y + 0.22, 0.52, 0.32, size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.82, y + 0.18, w - 1.0, 0.34, size=14, color=INK_LIGHT, bold=True)
    add_text(slide, body, x + 0.18, y + 0.78, w - 0.36, h - 0.95, size=12, color=rgb("C8D3E2"))


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "SOLANA HACKATHON 2026", 0.7, 0.78, w=2.55, fill=GOLD)
    add_text(slide, "Agent1c.ai", 0.7, 1.45, 5.4, 0.9, size=32, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_text(
        slide,
        "The browser-native agentic desktop for wallet-native work.",
        0.7,
        2.2,
        6.2,
        0.6,
        size=22,
        color=MINT,
        bold=True,
        font=TITLE_FONT,
    )
    add_text(
        slide,
        "Hitomi lives inside a real workspace, authenticates with Solana, and explains live wallet state in plain English instead of making users bounce between tabs, explorers, and dashboards.",
        0.7,
        3.0,
        5.6,
        1.25,
        size=16,
        color=rgb("C6D0DE"),
    )
    add_rect(slide, 0.7, 5.15, 5.75, 1.2, PANEL_DARK, radius=True)
    for idx, label in enumerate([
        "Wallet-native identity",
        "Persistent AI workspace",
        "Read-only on-chain context",
    ]):
        chip_x = 0.95 + (idx * 1.84)
        add_rect(slide, chip_x, 5.48, 1.55, 0.45, BG_DARK_2, line=MINT, radius=True)
        add_text(slide, label, chip_x + 0.1, 5.61, 1.35, 0.2, size=10, color=INK_LIGHT, align=PP_ALIGN.CENTER)
    add_oval(slide, 8.55, 1.25, 3.35, 3.35, rgb("113A45"))
    add_oval(slide, 8.9, 1.6, 2.65, 2.65, rgb("1B4A55"))
    if HEDGEHOG_IMAGE.exists():
        slide.shapes.add_picture(str(HEDGEHOG_IMAGE), Inches(8.95), Inches(1.45), width=Inches(2.6))
    add_text(
        slide,
        "Pitch deck grounded in the shipped repo state, with a technical appendix for deeper judge follow-up.",
        7.15,
        5.45,
        4.7,
        0.7,
        size=13,
        color=rgb("C6D0DE"),
    )
    add_footer(slide, 1, dark=True)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide)
    add_kicker(slide, "PROBLEM", 0.7, 0.62, fill=CORAL, text_color=INK_LIGHT)
    add_text(slide, "AI still lives beside the workflow.\nWeb3 still lives across too many surfaces.", 0.7, 1.05, 6.0, 1.4, size=28, color=INK_DARK, bold=True, font=TITLE_FONT)
    add_bullets(
        slide,
        [
            "Chatbots open in a separate tab and rarely understand the user's real environment.",
            "Wallet context is trapped in browser extensions, explorers, dashboards, and DMs.",
            "Even simple questions like \"what hit my wallet?\" require manual context switching and interpretation.",
        ],
        0.7,
        2.75,
        5.4,
        2.4,
        size=17,
        color=INK_DARK,
    )
    add_rect(slide, 7.0, 1.25, 5.2, 4.95, PANEL_LIGHT, line=rgb("E4DDD2"), radius=True)
    add_card(slide, 7.35, 1.65, 4.5, 0.88, "Today", "Chat tab, wallet extension, block explorer, notes, and Telegram all hold different slices of state.", fill=PANEL_SOFT)
    add_card(slide, 7.35, 2.75, 4.5, 0.88, "What breaks", "Context does not accumulate, identity does not travel cleanly, and users become the manual orchestrator.", fill=PANEL_SOFT)
    add_card(slide, 7.35, 3.85, 4.5, 0.88, "What users want", "One place where the assistant, files, browser, and wallet context already know each other.", fill=PANEL_SOFT)
    add_rect(slide, 0.7, 6.1, 11.45, 0.62, INK_DARK, radius=True)
    add_text(
        slide,
        "The missing product is not another chat box or another dashboard. It is a persistent agent workspace.",
        1.0,
        6.27,
        10.9,
        0.26,
        size=15,
        color=INK_LIGHT,
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_footer(slide, 2)


def slide_need(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "WHAT USERS ACTUALLY NEED", 0.7, 0.72, w=3.1, fill=MINT)
    add_text(slide, "A workspace-native agent, not a disposable conversation.", 0.7, 1.15, 7.0, 0.7, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_text(slide, "Agent1c is built around four product truths that normal copilots skip.", 0.7, 1.92, 6.0, 0.45, size=14, color=rgb("C7D3E2"))
    cards = [
        ("Persistent context", "The agent lives inside the desktop, so windows, docs, threads, and state stay in the same room."),
        ("Wallet-native identity", "The wallet becomes authenticated workspace context instead of staying stuck in a wallet popup."),
        ("Visible tool use", "The user can see and direct the workspace instead of trusting hidden background magic."),
        ("Safe first step", "Read-only wallet understanding is the right wedge before any signing or execution surface."),
    ]
    positions = [(0.7, 2.65), (6.8, 2.65), (0.7, 4.55), (6.8, 4.55)]
    for (title, body), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 5.85, 1.48, title, body, fill=PANEL_DARK, title_color=INK_LIGHT, body_color=rgb("C7D3E2"))
    add_footer(slide, 3, dark=True)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide, alt=True)
    add_kicker(slide, "SOLUTION", 0.7, 0.62, fill=MINT)
    add_title_block(
        slide,
        "Agent1c.ai is a hosted agentic desktop OS.",
        "It runs as a static web desktop with Hitomi, draggable windows, browser/file surfaces, editable runtime docs, and cloud-authenticated AI orchestration.",
    )
    add_bullets(
        slide,
        [
            "Frontend is a static GitHub Pages app with vanilla JS modules and no build pipeline for the core runtime.",
            "Hitomi operates inside the workspace instead of being limited to one chat panel.",
            "The same runtime can combine chat, files, browser actions, notes, and authenticated wallet context.",
        ],
        0.7,
        2.75,
        5.2,
        2.1,
        size=16,
    )
    add_rect(slide, 6.4, 1.45, 5.65, 4.85, PANEL_LIGHT, line=rgb("DDE5E8"), radius=True)
    add_card(slide, 6.72, 1.8, 1.35, 1.0, "User", "asks in plain English", fill=rgb("FFF7EA"), title_color=INK_DARK, body_color=INK_MUTED)
    add_card(slide, 8.35, 1.8, 1.6, 1.0, "Hitomi", "agent runtime", fill=rgb("EFF9F6"), title_color=INK_DARK, body_color=INK_MUTED)
    add_card(slide, 10.2, 1.8, 1.45, 1.0, "Tools", "wallet, browser, files", fill=rgb("EEF5FF"), title_color=INK_DARK, body_color=INK_MUTED)
    add_text(slide, "→", 7.98, 2.1, 0.25, 0.2, size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 9.93, 2.1, 0.25, 0.2, size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 7.1, 3.35, 4.55, 1.2, "Visible workspace", "Windows, notes, browser routing, desktop surfaces, and app state stay user-visible instead of being hidden in a backend black box.", fill=PANEL_SOFT)
    add_card(slide, 7.1, 4.85, 4.55, 1.0, "Cloud convenience, optional privacy upgrades", "Managed cloud is default on `.ai`; relay and Tor paths already exist as adjacent surfaces.", fill=PANEL_SOFT)
    add_footer(slide, 4)


def slide_solana(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "WHY SOLANA FITS", 0.7, 0.72, fill=GOLD)
    add_text(slide, "Solana is not a bolt-on feature here.\nIt is the identity wedge.", 0.7, 1.1, 6.5, 1.15, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_bullets(
        slide,
        [
            "The auth UI already includes `Continue with Solana` and uses Supabase `signInWithWeb3({ chain: \"solana\" })`.",
            "The runtime extracts the canonical wallet address from auth identity metadata instead of scraping labels.",
            "Hitomi can refresh balance and recent transactions through explicit read-only tools and answer in plain English.",
            "Longer term, the wallet becomes a clean anchor for encrypted continuity across devices and sessions.",
        ],
        0.7,
        2.65,
        5.85,
        3.0,
        size=16,
        color=rgb("D0DBE8"),
    )
    add_card(slide, 7.1, 1.65, 4.65, 1.05, "1. Sign in", "Phantom / Backpack / Brave / compatible Solana wallet establishes authenticated session.", fill=PANEL_DARK, title_color=MINT, body_color=rgb("C8D3E2"))
    add_card(slide, 7.1, 2.95, 4.65, 1.05, "2. Inspect", "Wallet state is cached in runtime and refreshed from Solana RPC with endpoint fallbacks.", fill=PANEL_DARK, title_color=SKY, body_color=rgb("C8D3E2"))
    add_card(slide, 7.1, 4.25, 4.65, 1.05, "3. Explain", "Hitomi answers questions like \"what's my balance?\" or \"what hit my wallet recently?\" safely and read-only.", fill=PANEL_DARK, title_color=GOLD, body_color=rgb("C8D3E2"))
    add_footer(slide, 5, dark=True)


def slide_walkthrough(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide)
    add_kicker(slide, "PRODUCT FLOW", 0.7, 0.62, fill=CORAL, text_color=INK_LIGHT)
    add_title_block(
        slide,
        "What a judge or user experiences in the product",
        "The story is simple enough for a pitch demo, but it sits on top of real runtime plumbing already in the repo.",
    )
    add_number_card(slide, 1, "Connect wallet", "Use Solana sign-in from the auth window and land in an authenticated cloud workspace.", 0.7, 2.45, 2.8, 1.55, fill=INK_DARK, accent=MINT)
    add_number_card(slide, 2, "Open the desktop", "Hitomi, notes, docs, credits, browser, and other panels all exist in the same visible workspace.", 3.62, 2.45, 2.8, 1.55, fill=INK_DARK, accent=SKY)
    add_number_card(slide, 3, "Ask naturally", "Questions like \"what's my balance?\" or \"show recent transactions\" stay in plain English.", 6.54, 2.45, 2.8, 1.55, fill=INK_DARK, accent=CORAL)
    add_number_card(slide, 4, "Get live context", "Hitomi calls the wallet overview or refresh tool, then answers from current wallet state.", 9.46, 2.45, 2.8, 1.55, fill=INK_DARK, accent=GOLD)
    add_rect(slide, 0.7, 4.55, 11.55, 1.35, PANEL_LIGHT, line=rgb("E3DACE"), radius=True)
    add_text(
        slide,
        "Important detail: the Solana scope is intentionally read-only today.",
        1.0,
        4.82,
        4.8,
        0.26,
        size=15,
        color=INK_DARK,
        bold=True,
    )
    add_text(
        slide,
        "That makes the demo safer, easier to trust, and easier to explain: Agent1c understands the wallet before it ever asks the user to approve anything.",
        1.0,
        5.18,
        10.7,
        0.42,
        size=14,
        color=INK_MUTED,
    )
    add_footer(slide, 6)


def slide_shipped(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "PROOF IT EXISTS", 0.7, 0.72, fill=MINT)
    add_text(slide, "This is already in the repo and in `origin/master`.", 0.7, 1.08, 6.8, 0.72, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_text(slide, "We are not asking the judges to fund a mockup. The core Solana path is already wired.", 0.7, 1.88, 6.2, 0.4, size=14, color=rgb("C7D3E2"))
    cards = [
        ("`js/agent1cauth.js`", "Wallet detection plus Supabase `signInWithWeb3` for Solana login."),
        ("`js/solana-wallet.js`", "Live RPC balance reads, recent signatures, and parsed transaction summaries."),
        ("`js/agent1c.js`", "Canonical wallet state bucket plus post-auth refresh orchestration."),
        ("`js/agent1c-tools-runtime.js`", "`solana_wallet_overview` and `solana_wallet_refresh` for Hitomi."),
        ("`js/agent1c-instance-profiles.js`", "Shipped runtime instructions that keep wallet behavior safe and read-only."),
    ]
    positions = [(0.7, 2.6), (6.6, 2.6), (0.7, 4.15), (6.6, 4.15), (3.65, 5.7)]
    widths = [5.2, 5.2, 5.2, 5.2, 5.95]
    for (title, body), (x, y), width in zip(cards, positions, widths):
        add_card(slide, x, y, width, 1.18, title, body, fill=PANEL_DARK, title_color=MINT, body_color=rgb("C8D3E2"))
    add_footer(slide, 7, dark=True)


def slide_moat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide, alt=True)
    add_kicker(slide, "WHY IT WINS", 0.7, 0.62, fill=MINT)
    add_title_block(
        slide,
        "Agent1c is not another chatbot and not another wallet dashboard.",
        "The product edge comes from combining workspace-native AI, user-visible tools, and wallet-native context in one surface.",
    )
    moat_cards = [
        ("Workspace-native", "The agent lives inside a desktop shell with windows, notes, files, browser, and panel state."),
        ("Wallet-aware", "Identity and wallet context can enter the workspace through authenticated Solana login."),
        ("Tool-using", "Hitomi does not have to guess. She can call explicit wallet tools and other runtime tools."),
        ("Privacy-flexible", "Cloud first on `.ai`, with relay and Tor surfaces already present for privacy upgrades."),
        ("Safe by default", "Current chain scope is read-only, which keeps claims and trust boundaries crisp."),
    ]
    y = 2.5
    for idx, (title, body) in enumerate(moat_cards):
        add_card(slide, 0.7 + (idx % 2) * 6.1, y + (idx // 2) * 1.35, 5.45, 1.08, title, body, fill=PANEL_LIGHT)
    add_rect(slide, 0.7, 6.28, 11.7, 0.52, INK_DARK, radius=True)
    add_text(
        slide,
        "The wedge is simple: turn wallet identity into working context for an always-available AI teammate.",
        0.95,
        6.42,
        11.2,
        0.2,
        size=14,
        color=INK_LIGHT,
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_footer(slide, 8)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "ARCHITECTURE", 0.7, 0.72, fill=GOLD)
    add_text(slide, "A simple product story on top of a modular runtime.", 0.7, 1.08, 6.8, 0.7, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_card(slide, 0.7, 2.1, 3.2, 1.55, "Frontend shell", "Static GitHub Pages app, HedgeyOS desktop shell, no npm build step for the core web app.", fill=PANEL_DARK, title_color=MINT, body_color=rgb("C8D3E2"))
    add_card(slide, 4.2, 2.1, 4.0, 1.55, "Agent runtime", "`js/agent1c.js` plus tool/runtime modules manage Hitomi, threads, docs, events, and wallet state.", fill=PANEL_DARK, title_color=SKY, body_color=rgb("C8D3E2"))
    add_card(slide, 8.55, 2.1, 3.95, 1.55, "Cloud services", "Supabase Auth + Edge Function `xai-chat` provide authentication and managed model access.", fill=PANEL_DARK, title_color=GOLD, body_color=rgb("C8D3E2"))
    add_text(slide, "↓", 2.15, 3.8, 0.3, 0.2, size=20, color=MINT, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, "↓", 6.0, 3.8, 0.3, 0.2, size=20, color=SKY, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, "↓", 10.35, 3.8, 0.3, 0.2, size=20, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_card(slide, 0.95, 4.3, 3.3, 1.45, "Solana RPC read layer", "`js/solana-wallet.js` calls `getBalance`, `getSignaturesForAddress`, and `getTransaction` with RPC fallbacks.", fill=PANEL_DARK, title_color=MINT, body_color=rgb("C8D3E2"))
    add_card(slide, 4.55, 4.3, 3.3, 1.45, "Prompt + tools contract", "Shipped `SOUL`, `TOOLS`, and heartbeat templates tell Hitomi how to handle wallet state safely.", fill=PANEL_DARK, title_color=SKY, body_color=rgb("C8D3E2"))
    add_card(slide, 8.15, 4.3, 3.3, 1.45, "Relay surfaces", "Shell relay, Tor relay, and proxy browsing broaden what the desktop can do beyond wallet reads.", fill=PANEL_DARK, title_color=GOLD, body_color=rgb("C8D3E2"))
    add_footer(slide, 9, dark=True)


def slide_trust(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide)
    add_kicker(slide, "TRUST MODEL", 0.7, 0.62, fill=CORAL, text_color=INK_LIGHT)
    add_title_block(
        slide,
        "What the Solana feature does, and what it very intentionally does not do.",
        "This slide matters because trust boundaries are part of the product, not just part of compliance.",
    )
    add_card(slide, 0.7, 2.4, 5.65, 3.5, "Does today", "", fill=PANEL_LIGHT)
    add_bullets(
        slide,
        [
            "Uses the authenticated wallet address as the canonical source of wallet identity.",
            "Fetches SOL balance and recent transactions through explicit read-only tool calls.",
            "Lets Hitomi answer wallet questions from refreshed runtime state instead of guessing.",
            "Keeps user-visible workspace surfaces in front of the user while the agent works.",
        ],
        1.0,
        2.95,
        5.05,
        2.6,
        size=15,
    )
    add_card(slide, 6.95, 2.4, 5.35, 3.5, "Does not do today", "", fill=PANEL_LIGHT)
    add_bullets(
        slide,
        [
            "No signing, sending, approving, swapping, or custody claims.",
            "No label scraping as primary production logic for wallet identity.",
            "No aggressive background polling loop pretending to keep wallet state magically fresh.",
            "No claim that Agent1c controls private keys or can move funds.",
        ],
        7.25,
        2.95,
        4.7,
        2.6,
        size=15,
    )
    add_footer(slide, 10)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "ROADMAP", 0.7, 0.72, fill=MINT)
    add_text(slide, "The wedge is read-only wallet understanding.\nThe platform ambition is much bigger.", 0.7, 1.08, 7.0, 1.0, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    phases = [
        ("Now", "Shipped", "Solana login, wallet identity extraction, balance + recent transaction answers."),
        ("Next", "Hackathon polish", "Tighter demo flow, better wallet summary surfaces, richer narrative for judges."),
        ("Then", "Product expansion", "SPL token views, deeper transaction explanations, wallet summary window."),
        ("Later", "Platform payoff", "Wallet-anchored encrypted sync and carefully approved action-taking surfaces."),
    ]
    for idx, (when, title, body) in enumerate(phases):
        x = 0.7 + idx * 3.1
        accent = [MINT, SKY, CORAL, GOLD][idx]
        add_rect(slide, x, 2.65, 2.75, 2.65, PANEL_DARK, radius=True)
        add_rect(slide, x, 2.65, 2.75, 0.16, accent)
        add_text(slide, when, x + 0.18, 2.95, 1.0, 0.24, size=12, color=accent, bold=True)
        add_text(slide, title, x + 0.18, 3.28, 2.35, 0.34, size=15, color=INK_LIGHT, bold=True)
        add_text(slide, body, x + 0.18, 3.78, 2.35, 1.2, size=12, color=rgb("C8D3E2"))
    add_footer(slide, 11, dark=True)


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide, alt=True)
    add_kicker(slide, "CLOSING", 0.7, 0.62, fill=MINT)
    add_text(slide, "Wallets should anchor AI workspaces,\nnot just transactions.", 0.7, 1.12, 6.7, 1.2, size=30, color=INK_DARK, bold=True, font=TITLE_FONT)
    add_bullets(
        slide,
        [
            "Agent1c makes Solana identity useful inside a living browser workspace.",
            "The repo already proves the core wallet-aware experience is real, not hypothetical.",
            "This hackathon is the right stage to turn that wedge into a memorable product demo.",
        ],
        0.7,
        2.85,
        5.7,
        2.2,
        size=18,
    )
    add_rect(slide, 7.2, 1.4, 4.6, 4.8, PANEL_LIGHT, line=rgb("DDE5E8"), radius=True)
    add_text(slide, "Agent1c.ai", 7.7, 1.95, 3.6, 0.45, size=24, color=INK_DARK, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, "Hosted agentic desktop OS", 7.7, 2.45, 3.6, 0.25, size=14, color=INK_MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.1, 3.0, 2.8, 0.55, INK_DARK, radius=True)
    add_text(slide, "Solana login + live wallet context", 8.22, 3.18, 2.56, 0.2, size=12, color=INK_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.1, 3.75, 2.8, 0.55, INK_DARK, radius=True)
    add_text(slide, "Hitomi answers from tools, not guesses", 8.22, 3.93, 2.56, 0.2, size=12, color=INK_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.1, 4.5, 2.8, 0.55, INK_DARK, radius=True)
    add_text(slide, "Built for a bigger workspace future", 8.22, 4.68, 2.56, 0.2, size=12, color=INK_LIGHT, align=PP_ALIGN.CENTER)
    if HEDGEHOG_IMAGE.exists():
        slide.shapes.add_picture(str(HEDGEHOG_IMAGE), Inches(9.0), Inches(5.15), width=Inches(1.1))
    add_footer(slide, 12)


def slide_appendix_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "APPENDIX", 0.7, 0.72, fill=SKY, text_color=BG_DARK)
    add_text(slide, "Technical stack snapshot", 0.7, 1.08, 5.4, 0.6, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    stack_cards = [
        ("Frontend", "Static GitHub Pages app\nVanilla JS modules\nHedgeyOS desktop shell\nNo build step for core app"),
        ("Core runtime", "`js/main.js`\n`js/wm.js`\n`js/agent1c.js`\n`js/agent1cauth.js`\n`js/agent1c-tools-runtime.js`"),
        ("Cloud layer", "Supabase project `gkfhxhrleuauhnuewfmw`\nSupabase Auth\nEdge Function `xai-chat`\nManaged usage/quota path"),
        ("Solana read layer", "`js/solana-wallet.js`\nRPC fallbacks\nBalance + signatures + tx parsing\nRead-only by design"),
    ]
    positions = [(0.7, 2.2), (3.95, 2.2), (7.2, 2.2), (10.45, 2.2)]
    for (title, body), (x, y) in zip(stack_cards, positions):
        add_card(slide, x, y, 2.55, 3.05, title, body, fill=PANEL_DARK, title_color=MINT, body_color=rgb("C8D3E2"))
    add_footer(slide, 13, appendix=True, dark=True)


def slide_appendix_auth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide)
    add_kicker(slide, "APPENDIX", 0.7, 0.62, fill=SKY)
    add_title_block(
        slide,
        "Solana auth pipeline",
        "This is the exact shape the repo describes and implements today.",
    )
    nodes = [
        ("Wallet provider", "Phantom / Backpack / Brave / compatible wallet"),
        ("`detectSolanaWallet()`", "Auth UI finds a connected provider in the browser"),
        ("Supabase `signInWithWeb3`", "Authenticated cloud session established with `chain: \"solana\"`"),
        ("Identity parser", "Wallet address normalized from identity metadata and custom claims"),
        ("Runtime state", "`applyCloudIdentityToWalletState()` stores canonical wallet context"),
    ]
    x_positions = [0.7, 3.05, 5.4, 7.9, 10.35]
    for idx, ((title, body), x) in enumerate(zip(nodes, x_positions)):
        fill = [PANEL_SOFT, PANEL_LIGHT, PANEL_SOFT, PANEL_LIGHT, PANEL_SOFT][idx]
        add_card(slide, x, 2.95, 2.1, 2.0, title, body, fill=fill)
        if idx < len(nodes) - 1:
            add_text(slide, "→", x + 2.12, 3.7, 0.18, 0.2, size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 14, appendix=True)


def slide_appendix_wallet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "APPENDIX", 0.7, 0.72, fill=SKY, text_color=BG_DARK)
    add_text(slide, "How a wallet question becomes an answer", 0.7, 1.08, 7.2, 0.65, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    steps = [
        ("1. User asks", "Example: \"What's in my wallet?\""),
        ("2. Prompt context", "Wallet address is inserted into runtime context if authenticated"),
        ("3. Tool call", "`solana_wallet_overview` or `solana_wallet_refresh`"),
        ("4. Runtime refresh", "`refreshConnectedWalletState()` updates wallet cache"),
        ("5. RPC reads", "`getBalance`, `getSignaturesForAddress`, `getTransaction`"),
        ("6. Plain-English answer", "Hitomi responds from tool output, not from guessing"),
    ]
    y = 2.15
    for idx, (title, body) in enumerate(steps):
        x = 0.7 + (idx % 3) * 4.1
        row = idx // 3
        add_card(slide, x, y + row * 1.75, 3.55, 1.25, title, body, fill=PANEL_DARK, title_color=GOLD if row == 0 else MINT, body_color=rgb("C8D3E2"))
    add_footer(slide, 15, appendix=True, dark=True)


def slide_appendix_breadth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_light(slide, alt=True)
    add_kicker(slide, "APPENDIX", 0.7, 0.62, fill=SKY)
    add_title_block(
        slide,
        "The Solana wedge sits inside a broader platform",
        "These adjacent surfaces matter because they show Agent1c is a workspace product, not a one-feature wallet demo.",
    )
    cards = [
        ("Browser + relay", "The repo already includes shell relay, Tor relay, and proxy browsing surfaces."),
        ("Telegram cloud relay", "Agent1c can bridge Telegram into the runtime while keeping tab-online constraints explicit."),
        ("Editable runtime docs", "`SOUL.md`, `TOOLS.md`, and heartbeat are first-class windows that shape behavior."),
        ("Managed cloud quotas", "Per-user usage is enforced in Supabase, not by trusting provider-global usage."),
    ]
    positions = [(0.7, 2.55), (6.65, 2.55), (0.7, 4.2), (6.65, 4.2)]
    for (title, body), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 5.0, 1.28, title, body, fill=PANEL_LIGHT)
    add_footer(slide, 16, appendix=True)


def slide_appendix_placeholders(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_dark(slide)
    add_kicker(slide, "APPENDIX", 0.7, 0.72, fill=SKY, text_color=BG_DARK)
    add_text(slide, "Reserved blocks to fill before finals", 0.7, 1.08, 6.8, 0.62, size=28, color=INK_LIGHT, bold=True, font=TITLE_FONT)
    add_text(slide, "These are intentionally structured into the deck now so we can add external research later without redesigning the narrative.", 0.7, 1.82, 7.2, 0.4, size=14, color=rgb("C7D3E2"))
    placeholders = [
        ("Market size", "TAM / SAM / SOM or the specific market lens we want to use."),
        ("ICP + wedge", "Which user segment we pitch first: power users, traders, teams, or creators."),
        ("Traction + proof", "Usage stats, qualitative demos, retention, waitlist, or judge-relevant proof points."),
        ("Competition", "How we position Agent1c against copilots, wallet dashboards, and browser assistants."),
    ]
    positions = [(0.7, 2.65), (6.6, 2.65), (0.7, 4.45), (6.6, 4.45)]
    for (title, body), (x, y) in zip(placeholders, positions):
        add_card(slide, x, y, 5.1, 1.35, title, body, fill=PANEL_DARK, title_color=SKY, body_color=rgb("C8D3E2"))
    add_footer(slide, 17, appendix=True, dark=True)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_problem(prs)
    slide_need(prs)
    slide_solution(prs)
    slide_solana(prs)
    slide_walkthrough(prs)
    slide_shipped(prs)
    slide_moat(prs)
    slide_architecture(prs)
    slide_trust(prs)
    slide_roadmap(prs)
    slide_close(prs)
    slide_appendix_stack(prs)
    slide_appendix_auth(prs)
    slide_appendix_wallet(prs)
    slide_appendix_breadth(prs)
    slide_appendix_placeholders(prs)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build_deck()
    print(OUT_PATH)
